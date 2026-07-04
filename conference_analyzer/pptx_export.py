"""Render an :class:`AnalysisResult` to a PowerPoint slide deck.

Generated deterministically with ``python-pptx`` (no API call). The styling
mirrors the web UI: white slides, a single navy accent, clean typography.

Deck structure:

* a **title** slide (theme, source, headline counts);
* an **overview** slide with a papers-per-topic bar chart; and
* one or more **topic** slides, each listing its papers with author lines and
  clickable PDF links (continuation slides are added when a topic overflows).
"""

from __future__ import annotations

import io
from typing import TYPE_CHECKING

if TYPE_CHECKING:  # avoid importing the model at module load
    from .models import AnalysisResult

# Palette shared with the UI -------------------------------------------------
_NAVY = (0x1F, 0x4E, 0x79)
_ACCENT = (0x2B, 0x6C, 0xB0)
_INK = (0x1A, 0x20, 0x2C)
_MUTED = (0x64, 0x74, 0x8B)
_LINE = (0xE2, 0xE8, 0xF0)
_WHITE = (0xFF, 0xFF, 0xFF)

_PAPERS_PER_SLIDE = 8


def build_pptx(result: "AnalysisResult") -> bytes:
    """Build the deck and return it as ``.pptx`` bytes.

    Raises ``RuntimeError`` with an actionable message if ``python-pptx`` is
    not installed.
    """
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Inches, Pt
    except ImportError as e:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "PPTX export needs the 'python-pptx' package. Install it with "
            "`pip install python-pptx`."
        ) from e

    def rgb(c: tuple[int, int, int]) -> "RGBColor":
        return RGBColor(*c)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    # ---- small helpers ------------------------------------------------ #
    def add_slide():
        return prs.slides.add_slide(blank)

    def fill_rect(slide, left, top, width, height, color):
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def textbox(slide, left, top, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        return box, tf

    def style_run(run, text, size, color, bold=False, italic=False):
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color)
        run.font.name = "Calibri"

    def header_bar(slide, title, subtitle=""):
        fill_rect(slide, 0, 0, SW, Inches(1.0), _NAVY)
        _, tf = textbox(slide, Inches(0.5), Inches(0.12), SW - Inches(1.0), Inches(0.8))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        style_run(p.add_run(), title, 24, _WHITE, bold=True)
        if subtitle:
            p2 = tf.add_paragraph()
            style_run(p2.add_run(), subtitle, 12, _WHITE)

    # ---- 1. title slide ---------------------------------------------- #
    s = add_slide()
    fill_rect(s, 0, 0, SW, SH, _WHITE)
    fill_rect(s, 0, Inches(2.4), SW, Inches(0.06), _ACCENT)

    _, tf = textbox(s, Inches(0.9), Inches(1.2), SW - Inches(1.8), Inches(1.2))
    style_run(tf.paragraphs[0].add_run(), "Conference Paper Analysis", 20, _MUTED)
    p = tf.add_paragraph()
    style_run(p.add_run(), result.theme, 40, _NAVY, bold=True)

    _, tf = textbox(s, Inches(0.9), Inches(2.7), SW - Inches(1.8), Inches(3.0))
    line = tf.paragraphs[0]
    style_run(line.add_run(), f"Source: {result.event_url}", 14, _MUTED)
    stats = [
        (str(result.scanned), "papers scanned"),
        (str(len(result.relevant_papers)), f"match “{result.theme}”"),
        (str(len(result.topics)), "topics discovered"),
    ]
    for value, label in stats:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        style_run(p.add_run(), f"{value}  ", 28, _ACCENT, bold=True)
        style_run(p.add_run(), label, 16, _INK)

    # ---- 2. overview / chart ----------------------------------------- #
    if result.topics:
        s = add_slide()
        fill_rect(s, 0, 0, SW, SH, _WHITE)
        header_bar(s, "Papers per topic", result.theme)
        chart_data = CategoryChartData()
        # reversed so the largest topic sits at the top of a horizontal bar chart
        topics_sorted = list(result.topics)
        chart_data.categories = [t.name for t in reversed(topics_sorted)]
        chart_data.add_series("Papers", [t.count for t in reversed(topics_sorted)])
        gframe = s.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.5), Inches(1.2), SW - Inches(1.0), SH - Inches(1.6),
            chart_data,
        )
        chart = gframe.chart
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(11)
        plot.data_labels.font.color.rgb = rgb(_INK)
        plot.gap_width = 60
        series = chart.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = rgb(_ACCENT)
        for axis in (chart.category_axis, chart.value_axis):
            axis.tick_labels.font.size = Pt(11)
            axis.tick_labels.font.color.rgb = rgb(_INK)

    # ---- 3. one (or more) slide(s) per topic ------------------------- #
    by_id = {p.paper_id: p for p in result.relevant_papers}
    for topic in result.topics:
        papers = [by_id[pid] for pid in topic.paper_ids if pid in by_id]
        papers.sort(key=lambda p: p.confidence or 0, reverse=True)

        # 3a. topic overview slide: description + common findings
        if topic.description or topic.findings:
            s = add_slide()
            fill_rect(s, 0, 0, SW, SH, _WHITE)
            header_bar(
                s,
                topic.name,
                f"{topic.count} paper{'s' if topic.count != 1 else ''}",
            )
            top = Inches(1.25)
            if topic.description:
                _, tf = textbox(s, Inches(0.5), top, SW - Inches(1.0), Inches(0.9))
                style_run(tf.paragraphs[0].add_run(), topic.description, 14, _INK)
                top = Inches(2.05)
            if topic.findings:
                _, tf = textbox(s, Inches(0.5), top, SW - Inches(1.0), SH - top - Inches(0.4))
                head = tf.paragraphs[0]
                style_run(head.add_run(), "Main findings across this topic", 12, _ACCENT, bold=True)
                for finding in topic.findings[:10]:
                    fp = tf.add_paragraph()
                    fp.space_before = Pt(5)
                    style_run(fp.add_run(), "•  " + finding, 13, _INK)
        chunks = [
            papers[i : i + _PAPERS_PER_SLIDE]
            for i in range(0, max(len(papers), 1), _PAPERS_PER_SLIDE)
        ] or [[]]
        for idx, chunk in enumerate(chunks):
            s = add_slide()
            fill_rect(s, 0, 0, SW, SH, _WHITE)
            suffix = "" if len(chunks) == 1 else f"  ({idx + 1}/{len(chunks)})"
            header_bar(
                s,
                f"{topic.name} — Papers{suffix}",
                f"{topic.count} paper{'s' if topic.count != 1 else ''}",
            )
            top = Inches(1.25)
            _, tf = textbox(s, Inches(0.5), top, SW - Inches(1.0), SH - top - Inches(0.3))
            first = True
            for paper in chunk:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_after = Pt(8)
                # title (clickable -> landing page when available)
                trun = p.add_run()
                style_run(trun, "• " + paper.title, 14, _NAVY, bold=True)
                if paper.url:
                    trun.hyperlink.address = paper.url
                if paper.confidence is not None:
                    style_run(p.add_run(), f"   ({paper.confidence:.0%})", 11, _MUTED)
                # PDF link (only when a PDF exists)
                if paper.pdf_url:
                    style_run(p.add_run(), "   ", 14, _INK)
                    pdf = p.add_run()
                    style_run(pdf, "[PDF]", 12, _ACCENT, bold=True)
                    pdf.hyperlink.address = paper.pdf_url
                # authors on a second line
                if paper.authors:
                    authors = ", ".join(paper.authors[:6]) + (
                        "…" if len(paper.authors) > 6 else ""
                    )
                    a = tf.add_paragraph()
                    a.space_after = Pt(8)
                    style_run(a.add_run(), authors, 11, _MUTED)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()
